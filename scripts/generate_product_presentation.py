from __future__ import annotations

import sqlite3
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
VENDOR = ROOT / ".codex_vendor"

if VENDOR.exists():
    sys.path.insert(0, str(VENDOR))

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt
except ModuleNotFoundError as exc:
    raise SystemExit(
        "Missing dependency: install python-pptx and pillow to regenerate this deck."
    ) from exc


SLIDE_W = 13.333
SLIDE_H = 7.5

COLORS = {
    "bg": RGBColor(0xF5, 0xF0, 0xE8),
    "bg_alt": RGBColor(0xEF, 0xE5, 0xD1),
    "panel": RGBColor(0xFF, 0xFA, 0xF0),
    "ink": RGBColor(0x1D, 0x24, 0x33),
    "muted": RGBColor(0x5D, 0x66, 0x78),
    "line": RGBColor(0xD6, 0xCC, 0xBB),
    "buy": RGBColor(0x0F, 0x8B, 0x6F),
    "hold": RGBColor(0xD4, 0x8C, 0x14),
    "sell": RGBColor(0xB7, 0x47, 0x3A),
    "accent": RGBColor(0x0C, 0x50, 0x6D),
    "accent_dark": RGBColor(0x10, 0x22, 0x2F),
    "accent_soft": RGBColor(0xDC, 0xEE, 0xF5),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

FONTS = {
    "display": "Avenir Next",
    "body": "Palatino Linotype",
}


def load_metrics() -> dict[str, int]:
    metrics = {
        "sources": 5,
        "core_assets": 8,
        "narrative_clusters": 4,
        "pipeline_runs": 0,
        "eval_runs": 0,
        "tweets": 0,
        "decision_reviews": 0,
        "proposed_reviews": 0,
        "approved_reviews": 0,
        "runtime_jobs": 0,
        "notifications": 0,
    }
    db_path = ROOT / "data" / "x-ticker.sqlite"

    if not db_path.exists():
        return metrics

    with sqlite3.connect(db_path) as conn:
        cur = conn.cursor()
        queries = {
            "pipeline_runs": "select count(*) from pipeline_runs",
            "eval_runs": "select count(*) from eval_runs",
            "tweets": "select count(*) from tweets",
            "decision_reviews": "select count(*) from decision_reviews",
            "approved_reviews": "select count(*) from decision_reviews where status='approved'",
            "proposed_reviews": "select count(*) from decision_reviews where status='proposed'",
            "runtime_jobs": "select count(*) from runtime_jobs",
            "notifications": "select count(*) from notification_events",
        }

        for label, query in queries.items():
            try:
                cur.execute(query)
                metrics[label] = int(cur.fetchone()[0] or 0)
            except sqlite3.Error:
                metrics[label] = 0

    return metrics


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def set_background(slide, dark: bool = False) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["accent_dark"] if dark else COLORS["bg"]

    shapes = slide.shapes
    if dark:
        blob_specs = [
            (9.9, -0.8, 3.0, 3.0, COLORS["accent"], 0.55),
            (-0.8, 5.5, 2.6, 2.6, COLORS["sell"], 0.72),
            (10.9, 5.7, 2.2, 2.2, COLORS["buy"], 0.68),
        ]
    else:
        blob_specs = [
            (10.4, -0.6, 3.2, 3.2, COLORS["accent_soft"], 0.28),
            (-0.7, 5.6, 2.5, 2.5, COLORS["sell"], 0.85),
            (11.1, 5.9, 2.0, 2.0, COLORS["buy"], 0.88),
        ]

    for left, top, width, height, color, transparency in blob_specs:
        shape = shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.fill.transparency = transparency


def add_textbox(slide, left, top, width, height, text=""):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if text:
        box.text_frame.text = text
    return box


def style_title(shape, text: str, size: int, color: RGBColor, align=PP_ALIGN.LEFT) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = FONTS["display"]
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = color


def style_body(shape, lines: list[str], color: RGBColor, size: int = 19, bullet: bool = True) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = f"- {line}" if bullet else line
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_after = Pt(8)
        run = p.runs[0]
        run.font.name = FONTS["body"]
        run.font.size = Pt(size)
        run.font.color.rgb = color


def add_kicker(slide, text: str, left: float, top: float, dark: bool = False) -> None:
    box = add_textbox(slide, left, top, 4.3, 0.35)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text.upper()
    run = p.runs[0]
    run.font.name = FONTS["display"]
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = COLORS["bg_alt"] if dark else COLORS["accent"]


def add_footer(slide, text: str, dark: bool = False) -> None:
    box = add_textbox(slide, 0.6, 7.02, 12.0, 0.25)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.name = FONTS["display"]
    run.font.size = Pt(9)
    run.font.color.rgb = COLORS["bg_alt"] if dark else COLORS["muted"]


def add_tag(slide, text: str, left: float, top: float, fill_color: RGBColor, text_color: RGBColor) -> None:
    width = 0.48 + (0.072 * len(text))
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.34),
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.fill.transparency = 0.04

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = FONTS["display"]
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = text_color


def add_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    body_lines: list[str],
    accent_color: RGBColor,
    body_size: int = 16,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["panel"]
    shape.line.color.rgb = COLORS["line"]

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.12),
        Inches(height),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = accent_color
    accent.line.fill.background()

    title_box = add_textbox(slide, left + 0.28, top + 0.18, width - 0.42, 0.52)
    style_title(title_box, title, 19, COLORS["ink"])

    body_box = add_textbox(slide, left + 0.28, top + 0.72, width - 0.45, height - 0.92)
    style_body(body_box, body_lines, COLORS["muted"], size=body_size, bullet=False)


def add_stat_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    value: str,
    label: str,
    accent_color: RGBColor,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["panel"]
    shape.line.color.rgb = COLORS["line"]

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.08),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    value_box = add_textbox(slide, left + 0.18, top + 0.18, width - 0.3, 0.55)
    style_title(value_box, value, 22, accent_color)

    label_box = add_textbox(slide, left + 0.18, top + 0.74, width - 0.3, 0.38)
    tf = label_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    run = p.runs[0]
    run.font.name = FONTS["display"]
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = COLORS["muted"]


def add_section_divider(slide, title: str, subtitle: str) -> None:
    set_background(slide, dark=True)
    add_kicker(slide, "Section 02", 0.8, 0.8, dark=True)

    title_box = add_textbox(slide, 0.8, 1.3, 7.2, 1.5)
    style_title(title_box, title, 30, COLORS["white"])

    subtitle_box = add_textbox(slide, 0.8, 3.0, 5.3, 1.5)
    style_body(subtitle_box, [subtitle], COLORS["bg_alt"], size=19, bullet=False)

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.1),
        Inches(1.1),
        Inches(4.0),
        Inches(4.9),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["white"]
    panel.fill.transparency = 0.92
    panel.line.fill.background()

    steps = [
        ("1", "Interpret", COLORS["accent"]),
        ("2", "Constrain", COLORS["hold"]),
        ("3", "Approve", COLORS["buy"]),
    ]

    step_top = 1.7
    for number, label, color in steps:
        bubble = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(8.5),
            Inches(step_top),
            Inches(0.62),
            Inches(0.62),
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = color
        bubble.line.fill.background()
        tf = bubble.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = number
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = FONTS["display"]
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = COLORS["white"]

        step_box = add_textbox(slide, 9.3, step_top - 0.03, 2.3, 0.65)
        style_title(step_box, label, 18, COLORS["white"])
        step_top += 1.1

    add_footer(slide, "Sources: README.md, PRD.txt, src/pipelineRunner.js, src/modelClaimExtractor.js", dark=True)


def add_connector(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor) -> None:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(2.2)


def add_flow_box(slide, left: float, top: float, width: float, height: float, title: str, fill: RGBColor) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = FONTS["display"]
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"] if fill != COLORS["accent_soft"] else COLORS["accent_dark"]


def build_deck() -> Path:
    metrics = load_metrics()
    prs = new_presentation()
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_background(slide, dark=True)
    add_kicker(slide, "X Ticker Investment", 0.8, 0.65, dark=True)

    title_box = add_textbox(slide, 0.8, 1.05, 7.4, 1.8)
    style_title(title_box, "Research-first investment decisions from social signal noise", 28, COLORS["white"])

    subtitle_box = add_textbox(slide, 0.8, 2.9, 6.0, 1.1)
    style_body(
        subtitle_box,
        [
            "A local-first operator desk that turns curated X posts into explainable BUY / HOLD / SELL candidates, then forces research and approval before action."
        ],
        COLORS["bg_alt"],
        size=18,
        bullet=False,
    )

    add_tag(slide, "Bounded agentic pipeline", 0.8, 4.22, COLORS["accent"], COLORS["white"])
    add_tag(slide, "Human approval required", 3.3, 4.22, COLORS["hold"], COLORS["white"])
    add_tag(slide, "Replayable local runtime", 5.72, 4.22, COLORS["buy"], COLORS["white"])

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(8.45),
        Inches(0.9),
        Inches(3.95),
        Inches(5.6),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["white"]
    panel.fill.transparency = 0.9
    panel.line.fill.background()

    panel_title = add_textbox(slide, 8.8, 1.2, 2.7, 0.4)
    style_title(panel_title, "Current repo snapshot", 17, COLORS["white"])

    stat_specs = [
        (f"{metrics['sources']}", "Monitored sources", COLORS["accent"]),
        (f"{metrics['core_assets']}", "Core assets", COLORS["buy"]),
        (f"{metrics['pipeline_runs']}", "Pipeline runs", COLORS["hold"]),
        (f"{metrics['eval_runs']}", "Eval runs", COLORS["sell"]),
    ]
    stat_top = 1.8
    for value, label, color in stat_specs:
        add_stat_card(slide, 8.8, stat_top, 3.0, 1.0, value, label, color)
        stat_top += 1.08

    note_box = add_textbox(slide, 8.8, 6.0, 3.0, 0.55)
    style_body(note_box, ["Focused on AI and tech signals with a multi-day decision horizon."], COLORS["bg_alt"], size=12, bullet=False)
    add_footer(slide, "Sources: README.md, PRD.txt, src/data.js, data/x-ticker.sqlite", dark=True)

    slide = prs.slides.add_slide(blank)
    set_background(slide, dark=False)
    add_kicker(slide, "Business Impact", 0.7, 0.5)
    title_box = add_textbox(slide, 0.7, 0.82, 6.8, 0.82)
    style_title(title_box, "Why this product matters", 24, COLORS["ink"])

    intro_box = add_textbox(slide, 0.7, 1.62, 7.8, 0.55)
    style_body(
        intro_box,
        ["The product replaces fragmented feed reading with a governed research workflow that improves speed, consistency, and trust."],
        COLORS["muted"],
        size=16,
        bullet=False,
    )

    add_card(
        slide,
        0.7,
        2.2,
        3.85,
        1.7,
        "Signal overload becomes narrative triage",
        [
            f"{metrics['tweets']} persisted posts already sit behind the product's runtime.",
            "Repeated posts are clustered into narratives instead of being treated as isolated events.",
        ],
        COLORS["accent"],
    )
    add_card(
        slide,
        4.72,
        2.2,
        3.85,
        1.7,
        "Decision quality is forced upstream",
        [
            "Research dossiers need evidence, contradictions, citations, and operator validation before approval.",
            "That lowers the odds of acting on hype, rumor, or partial context.",
        ],
        COLORS["hold"],
    )
    add_card(
        slide,
        8.74,
        2.2,
        3.85,
        1.7,
        "Coverage becomes personal and scalable",
        [
            "The watched universe blends curated assets with holdings and watchlist names.",
            "One operator gets broader monitoring without turning the system into an autonomous trader.",
        ],
        COLORS["buy"],
    )

    process_title = add_textbox(slide, 0.7, 4.45, 4.0, 0.35)
    style_title(process_title, "Value chain", 16, COLORS["accent"])

    stages = [
        ("Capture", COLORS["accent"]),
        ("Validate", COLORS["hold"]),
        ("Approve", COLORS["buy"]),
        ("Monitor", COLORS["sell"]),
    ]
    stage_left = 0.7
    for label, color in stages:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(stage_left),
            Inches(4.9),
            Inches(1.8),
            Inches(0.68),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = FONTS["display"]
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = COLORS["white"]
        stage_left += 2.05

    narrative_box = add_textbox(slide, 0.7, 6.0, 7.2, 0.55)
    style_body(
        narrative_box,
        ["Net effect: faster interpretation, fewer premature calls, and a reusable audit trail that compounds operator knowledge over time."],
        COLORS["ink"],
        size=15,
        bullet=False,
    )
    add_footer(slide, "Sources: README.md, PRD.txt, src/researchStore.js, src/decisionReviewStore.js")

    slide = prs.slides.add_slide(blank)
    set_background(slide, dark=False)
    add_kicker(slide, "Business Impact", 0.7, 0.5)
    title_box = add_textbox(slide, 0.7, 0.82, 7.6, 0.82)
    style_title(title_box, "Operating leverage is already visible in the repo", 24, COLORS["ink"])

    add_stat_card(slide, 0.7, 1.75, 1.9, 1.18, str(metrics["pipeline_runs"]), "persisted pipeline runs", COLORS["accent"])
    add_stat_card(slide, 2.78, 1.75, 1.9, 1.18, str(metrics["eval_runs"]), "eval runs logged", COLORS["hold"])
    add_stat_card(slide, 4.86, 1.75, 1.9, 1.18, str(metrics["decision_reviews"]), "decision reviews", COLORS["buy"])
    add_stat_card(slide, 6.94, 1.75, 1.9, 1.18, str(metrics["runtime_jobs"]), "runtime jobs", COLORS["sell"])
    add_stat_card(slide, 9.02, 1.75, 1.9, 1.18, str(metrics["notifications"]), "notifications sent", COLORS["accent"])

    left_box = add_textbox(slide, 0.7, 3.25, 5.65, 2.7)
    style_body(
        left_box,
        [
            f"{metrics['proposed_reviews']} proposals and {metrics['approved_reviews']} approved review already show the approval queue working as intended.",
            "The scheduler, digest, and notification layer reduce the time spent checking whether anything important changed.",
            "Replayable runs make debugging, postmortems, and operator trust easier than black-box automation.",
        ],
        COLORS["ink"],
        size=17,
        bullet=True,
    )

    right_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(7.0),
        Inches(3.15),
        Inches(5.65),
        Inches(2.9),
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = COLORS["accent_dark"]
    right_panel.line.fill.background()

    right_title = add_textbox(slide, 7.3, 3.45, 5.0, 0.55)
    style_title(right_title, "Why this matters commercially", 20, COLORS["white"])

    right_body = add_textbox(slide, 7.3, 4.15, 4.9, 1.5)
    style_body(
        right_body,
        [
            "This is leverage through bounded automation, not autonomy.",
            "The product saves operator time while increasing explainability and lowering governance risk.",
            "That combination is the basis for premium decision-support positioning.",
        ],
        COLORS["bg_alt"],
        size=16,
        bullet=True,
    )
    add_footer(slide, "Sources: data/x-ticker.sqlite, src/pipelineRunner.js, src/orchestrator.js, src/reportBuilder.js")

    slide = prs.slides.add_slide(blank)
    add_section_divider(
        slide,
        "The AI agentic approach",
        "Specialist model steps interpret messy inputs, but deterministic rules, research gates, and operator review decide what becomes action-ready.",
    )

    slide = prs.slides.add_slide(blank)
    set_background(slide, dark=False)
    add_kicker(slide, "AI Agentic Approach", 0.7, 0.5)
    title_box = add_textbox(slide, 0.7, 0.82, 6.5, 0.82)
    style_title(title_box, "Bounded architecture from signal to decision", 24, COLORS["ink"])

    flow_boxes = [
        (0.72, 1.75, 1.68, 0.72, "Curated X +\nmanual intake", COLORS["accent"]),
        (2.62, 1.75, 1.55, 0.72, "Ingestion +\nnormalization", COLORS["hold"]),
        (4.35, 1.75, 1.68, 0.72, "Claim extraction\n(OpenAI or fallback)", COLORS["buy"]),
        (6.22, 1.75, 1.5, 0.72, "Narrative\nclustering", COLORS["accent"]),
        (7.92, 1.75, 1.76, 0.72, "Policy + veto\nlayer", COLORS["sell"]),
        (9.88, 1.75, 1.55, 0.72, "Decision\nmath", COLORS["hold"]),
    ]
    for left, top, width, height, label, color in flow_boxes:
        add_flow_box(slide, left, top, width, height, label, color)

    for idx in range(len(flow_boxes) - 1):
        x1 = flow_boxes[idx][0] + flow_boxes[idx][2]
        y1 = flow_boxes[idx][1] + 0.36
        x2 = flow_boxes[idx + 1][0]
        y2 = flow_boxes[idx + 1][1] + 0.36
        add_connector(slide, x1, y1, x2, y2, COLORS["muted"])

    add_flow_box(slide, 2.0, 3.55, 2.1, 0.78, "Research dossier\nvalidated?", COLORS["accent_soft"])
    add_flow_box(slide, 5.0, 3.55, 2.1, 0.78, "Approval queue", COLORS["accent_soft"])
    add_flow_box(slide, 8.0, 3.55, 2.5, 0.78, "Advisor / dashboard /\ndigest / history", COLORS["accent_soft"])

    add_connector(slide, 10.65, 2.47, 10.65, 3.55, COLORS["muted"])
    add_connector(slide, 4.1, 3.94, 5.0, 3.94, COLORS["muted"])
    add_connector(slide, 7.1, 3.94, 8.0, 3.94, COLORS["muted"])

    note_box = add_textbox(slide, 0.72, 5.2, 11.6, 1.0)
    style_body(
        note_box,
        [
            "Design principle: AI interprets and normalizes. Deterministic policy, research validation, and human approval decide whether the output deserves trust."
        ],
        COLORS["ink"],
        size=16,
        bullet=False,
    )
    add_footer(slide, "Sources: README.md, src/pipelineRunner.js, src/modelClaimExtractor.js, src/agenticEngine.js")

    slide = prs.slides.add_slide(blank)
    set_background(slide, dark=False)
    add_kicker(slide, "AI Agentic Approach", 0.7, 0.5)
    title_box = add_textbox(slide, 0.7, 0.82, 6.6, 0.82)
    style_title(title_box, "Each module has a narrow job", 24, COLORS["ink"])

    cards = [
        (0.7, 1.7, "Claim extraction", ["Classifies post type, direction, explicitness, actionability, and mapped assets."], COLORS["accent"]),
        (4.55, 1.7, "Impact mapping", ["Ranks likely direct, read-through, and second-order effects across the watched universe."], COLORS["buy"]),
        (8.4, 1.7, "Narrative clustering", ["Groups repeated claims into event-like narratives so the product reasons at the right level."], COLORS["hold"]),
        (0.7, 4.0, "Policy and veto layer", ["Downgrades rumor, weak evidence, and noisy headlines before they contaminate the queue."], COLORS["sell"]),
        (4.55, 4.0, "Decision math", ["Surfaces thesis probability, upside, downside, reward-to-risk, size band, and max-loss guardrail."], COLORS["accent"]),
        (8.4, 4.0, "Portfolio-aware advisor", ["Answers from the latest snapshot plus saved holdings and watchlist context while respecting governance state."], COLORS["buy"]),
    ]
    for left, top, title, body, color in cards:
        add_card(slide, left, top, 3.45, 1.9, title, body, color, body_size=15)
    add_footer(slide, "Sources: src/modelClaimExtractor.js, src/modelImpactMapper.js, src/agenticEngine.js, src/financialAdvisor.js")

    slide = prs.slides.add_slide(blank)
    set_background(slide, dark=False)
    add_kicker(slide, "AI Agentic Approach", 0.7, 0.5)
    title_box = add_textbox(slide, 0.7, 0.82, 7.2, 0.82)
    style_title(title_box, "Governance, evals, and staged rollout keep the system safe", 24, COLORS["ink"])

    lifecycle_title = add_textbox(slide, 0.7, 1.7, 5.0, 0.35)
    style_title(lifecycle_title, "Research lifecycle", 16, COLORS["accent"])

    lifecycle_steps = [
        ("Discovery", COLORS["muted"]),
        ("Candidate", COLORS["accent"]),
        ("Validated", COLORS["hold"]),
        ("Approved", COLORS["buy"]),
        ("Dismissed / expired / archived", COLORS["sell"]),
    ]
    left = 0.7
    top = 2.2
    widths = [1.35, 1.35, 1.35, 1.2, 2.35]
    for (label, color), width in zip(lifecycle_steps, widths):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(0.65),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.name = FONTS["display"]
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = COLORS["white"]
        left += width + 0.14

    left_box = add_textbox(slide, 0.7, 3.25, 6.3, 2.2)
    style_body(
        left_box,
        [
            "A dossier cannot be validated or approved without thesis, assets, horizon, supporting evidence, contradicting evidence, and citations.",
            "Decisions only enter the queue when linked research is validated or approved.",
            "The advisor falls back to Research more or watch-only language when governance is incomplete.",
        ],
        COLORS["ink"],
        size=16,
        bullet=True,
    )

    add_stat_card(slide, 7.65, 1.82, 2.15, 1.15, str(metrics["eval_runs"]), "eval runs", COLORS["accent"])
    add_stat_card(slide, 9.95, 1.82, 2.15, 1.15, str(metrics["pipeline_runs"]), "pipeline replays", COLORS["buy"])
    add_stat_card(slide, 7.65, 3.2, 2.15, 1.15, str(metrics["decision_reviews"]), "review records", COLORS["hold"])
    add_stat_card(slide, 9.95, 3.2, 2.15, 1.15, str(metrics["notifications"]), "notifications", COLORS["sell"])

    right_body = add_textbox(slide, 7.65, 4.7, 4.55, 1.25)
    style_body(
        right_body,
        [
            "Extractor and model eval harnesses support field-level scoring, scenario tests, JSON reliability checks, and safer model changes.",
            "Rollout stays staged: research desk first, assisted approval next, bounded production only after the eval bar is met.",
        ],
        COLORS["muted"],
        size=15,
        bullet=True,
    )
    add_footer(slide, "Sources: README.md, docs/delta-roadmap.md, src/evalHarness.js, src/modelEvalHarness.js, src/researchStore.js")

    slide = prs.slides.add_slide(blank)
    set_background(slide, dark=True)
    add_kicker(slide, "Close", 0.8, 0.65, dark=True)
    title_box = add_textbox(slide, 0.8, 1.0, 7.4, 1.4)
    style_title(title_box, "The differentiator is disciplined agentic design", 28, COLORS["white"])

    body_box = add_textbox(slide, 0.8, 2.55, 6.8, 1.85)
    style_body(
        body_box,
        [
            "Most AI investing tools over-index on autonomy. This product wins by making interpretation faster while keeping trust, governance, and auditability in front.",
            "That is a stronger foundation for operator adoption, premium positioning, and safe expansion over time.",
        ],
        COLORS["bg_alt"],
        size=18,
        bullet=False,
    )

    roadmap_title = add_textbox(slide, 8.0, 1.2, 3.8, 0.35)
    style_title(roadmap_title, "Rollout path", 17, COLORS["white"])

    roadmap = [
        ("Stage 1", "Research desk mode"),
        ("Stage 2", "Assisted approval"),
        ("Stage 3", "Bounded production"),
    ]
    top = 1.75
    for stage, label in roadmap:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(8.0),
            Inches(top),
            Inches(3.9),
            Inches(0.95),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["white"]
        card.fill.transparency = 0.9
        card.line.fill.background()

        label_box = add_textbox(slide, 8.28, top + 0.14, 3.1, 0.5)
        style_title(label_box, f"{stage}: {label}", 15, COLORS["accent_dark"])
        top += 1.2

    add_footer(slide, "Sources: README.md, docs/delta-roadmap.md", dark=True)

    output_path = ROOT / "presentations" / "x-ticker-investment-product-story.pptx"
    prs.save(output_path)
    return output_path


if __name__ == "__main__":
    path = build_deck()
    print(path)
